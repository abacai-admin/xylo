import streamlit as st
import pandas as pd
import plotly.express as px
import plotly.graph_objects as go
from io import BytesIO
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
import os
from typing import Dict, List, Optional, Tuple, Any
import numpy as np

# Initialize session state for selected slides and charts
if 'selected_slides' not in st.session_state:
    st.session_state.selected_slides = {}

if 'selected_charts' not in st.session_state:
    st.session_state.selected_charts = {}

if 'selected_cagr_charts' not in st.session_state:
    st.session_state.selected_cagr_charts = {}

def create_presentation():
    """
    Create a PowerPoint presentation from the selected slides and charts.
    Uses the pptx_generator to create proper chart slides.
    """
    from logic.pptx_generator import generate_presentation
    from io import BytesIO
    import tempfile
    import os
    
    # Prepare slides data for the generator
    slides_data = []
    selected_slides = st.session_state.get('selected_slides', {})
    selected_charts = st.session_state.get('selected_charts', {})
    selected_cagr_charts = st.session_state.get('selected_cagr_charts', {})
    
    # Process each slide
    for idx, slide_data in enumerate(st.session_state.get('slides', [])):
        slide_id = slide_data.get('id', f'slide_{idx}')
        
        # Only process selected slides
        if selected_slides.get(slide_id, True):
            # Create a copy of the slide data to modify
            slide_copy = slide_data.copy()
            
            # Convert chart data to DataFrame if it's a dictionary
            if 'chart_data' in slide_copy:
                if isinstance(slide_copy['chart_data'], dict):
                    slide_copy['chart_data'] = pd.DataFrame(slide_copy['chart_data'])
                # Ensure it's not an empty DataFrame
                elif isinstance(slide_copy['chart_data'], pd.DataFrame) and slide_copy['chart_data'].empty:
                    continue  # Skip slides with empty data
            else:
                continue  # Skip slides with no chart data
            
            # If this is a comparison slide, ensure we have both datasets
            if slide_copy.get('is_comparison'):
                if 'chart_data2' in slide_copy:
                    if isinstance(slide_copy['chart_data2'], dict):
                        slide_copy['chart_data2'] = pd.DataFrame(slide_copy['chart_data2'])
                    # Ensure second dataset is not empty
                    if isinstance(slide_copy['chart_data2'], pd.DataFrame) and slide_copy['chart_data2'].empty:
                        continue  # Skip if second dataset is empty
                else:
                    continue  # Skip if comparison is True but no second dataset
            
            # Add to slides data
            slides_data.append(slide_copy)
    
    # Create a temporary file for the presentation
    with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as tmp_file:
        temp_path = tmp_file.name
    
    try:
        # Generate the presentation using the pptx_generator
        output_path = generate_presentation(slides_data, temp_path)
        
        # Read the generated presentation into memory
        with open(output_path, 'rb') as f:
            pptx_bytes = BytesIO(f.read())
        
        pptx_bytes.seek(0)
        return pptx_bytes
        
    except Exception as e:
        st.error(f"Error generating presentation: {e}")
        # Create a minimal error presentation
        prs = Presentation()
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_slide.shapes.title.text = "Error Generating Presentation"
        subtitle = title_slide.placeholders[1]
        subtitle.text = f"An error occurred: {str(e)}"
        
        # Save to bytes
        error_bytes = BytesIO()
        prs.save(error_bytes)
        error_bytes.seek(0)
        return error_bytes
    
    finally:
        # Clean up the temporary file
        if os.path.exists(temp_path):
            try:
                os.unlink(temp_path)
            except:
                pass

def display_slide_preview(slide, idx):
    """Display a preview of a single slide with selection options"""
    slide_id = slide.get('id', f'slide_{idx}')
    slide_title = slide.get('title', f'Untitled Slide {idx + 1}')
    
    # Initialize selection state if not exists
    if slide_id not in st.session_state.selected_slides:
        st.session_state.selected_slides[slide_id] = True
    if f"{slide_id}_chart" not in st.session_state.selected_charts:
        st.session_state.selected_charts[f"{slide_id}_chart"] = True
    if f"{slide_id}_cagr" not in st.session_state.selected_cagr_charts:
        st.session_state.selected_cagr_charts[f"{slide_id}_cagr"] = False
    
    with st.expander(f"{idx + 1}. {slide_title}", expanded=True):
        # Slide selection checkbox
        col1, col2 = st.columns([1, 3])
        with col1:
            st.session_state.selected_slides[slide_id] = st.checkbox(
                "Include in Export", 
                value=st.session_state.selected_slides[slide_id],
                key=f"slide_select_{slide_id}",
                help="Include this slide in the exported PowerPoint"
            )
        
        with col2:
            st.markdown(f"**{slide_title}**")
        
        # Display slide content
        if slide.get('content'):
            st.markdown("**Content:**")
            for point in slide.get('content', []):
                st.markdown(f"- {point}")
        
        # Display chart preview if available
        if slide.get('chart_data') is not None:
            st.markdown("---")
            st.markdown("**Chart:**")
            
            # Chart selection checkbox
            chart_col1, chart_col2 = st.columns([1, 3])
            with chart_col1:
                st.session_state.selected_charts[f"{slide_id}_chart"] = st.checkbox(
                    "Include Chart in Export",
                    value=st.session_state.selected_charts.get(f"{slide_id}_chart", True),
                    key=f"chart_select_{slide_id}",
                    help="Include this chart in the exported PowerPoint"
                )
            
            with chart_col2:
                chart_type_display = slide.get('chart_type', 'Chart').replace('_', ' ').title()
                # Get selected columns for more descriptive title
                selected_cols = slide.get('selected_columns', [])
                if not selected_cols and isinstance(slide.get('chart_data'), pd.DataFrame):
                    # Fallback for older slide structures or if selected_columns is empty
                    df_temp = pd.DataFrame(slide.get('chart_data'))
                    if not df_temp.empty:
                        selected_cols = df_temp.columns.tolist()
                        if 'Year' in selected_cols: selected_cols.remove('Year')
                        if 'Date' in selected_cols: selected_cols.remove('Date')
                
                metrics_text = ", ".join(selected_cols) if selected_cols else "Selected Metrics"
                st.markdown(f"**Visualization Type:** {chart_type_display} for ({metrics_text})")
            
            # Display chart preview
            try:
                if isinstance(slide['chart_data'], pd.DataFrame):
                    df = slide['chart_data']
                else:
                    df = pd.DataFrame(slide['chart_data'])
                
                # st.dataframe(df, use_container_width=True) # Removed raw data display
                
                # Simple chart visualization based on type
                chart_type = slide.get('chart_type', '').lower()
                # Ensure 'Year' or 'Date' is the index if present for st.bar_chart, st.line_chart, st.area_chart
                index_col = None
                if 'Year' in df.columns:
                    index_col = 'Year'
                elif 'Date' in df.columns:
                    index_col = 'Date'
                
                display_df = df.copy()
                if index_col and index_col in display_df.columns:
                    try:
                        display_df = display_df.set_index(index_col)
                    except Exception as e:
                        st.warning(f"Could not set index for chart preview: {e}")
                
                # Filter out non-numeric columns for relevant charts if selected_cols is not specific enough
                if selected_cols:
                    cols_to_plot = [col for col in selected_cols if col in display_df.columns and pd.api.types.is_numeric_dtype(display_df[col])]
                    if not cols_to_plot: # If selected_cols are not numeric, fallback to all numeric
                        cols_to_plot = display_df.select_dtypes(include=np.number).columns.tolist()
                else:
                    cols_to_plot = display_df.select_dtypes(include=np.number).columns.tolist()
                
                if not cols_to_plot and not (chart_type == 'table' and not df.empty):
                    st.info("No numeric data available to plot for the selected chart type.")
                elif chart_type == 'table' and not df.empty:
                    st.table(df) # Show table using st.table for better formatting than st.dataframe
                elif chart_type == 'bar_chart' and not display_df.empty and cols_to_plot:
                    st.bar_chart(display_df[cols_to_plot])
                elif chart_type == 'line_chart' and not display_df.empty and cols_to_plot:
                    st.line_chart(display_df[cols_to_plot])
                elif chart_type == 'pie' and not df.empty: # Pie chart uses original df structure typically
                    # Pie charts usually take a label column and a value column.
                    # Assuming first column is labels, second is values if not specified by selected_cols
                    if len(df.columns) >= 2:
                        label_col = df.columns[0]
                        value_col = df.columns[1]
                        if selected_cols and len(selected_cols) >= 2:
                            label_col = selected_cols[0] if selected_cols[0] in df.columns else df.columns[0]
                            value_col = selected_cols[1] if selected_cols[1] in df.columns else df.columns[1]
                        
                        if value_col in df.columns and pd.api.types.is_numeric_dtype(df[value_col]):
                            fig = px.pie(df, values=value_col, names=label_col, title=f"Pie Chart for {value_col}")
                            st.plotly_chart(fig, use_container_width=True)
                        else:
                            st.warning("Pie chart requires a numeric value column.")
                    else:
                        st.warning("Pie chart requires at least two columns (labels and values).")
                elif chart_type == 'area_chart' and not display_df.empty and cols_to_plot:
                    st.area_chart(display_df[cols_to_plot])
                elif not df.empty:
                    st.info(f"Preview for '{chart_type_display}' not fully implemented yet. Raw data shown as table.")
                    st.table(df) # Fallback to table if chart type specific preview isn't there
                
            except Exception as e:
                st.error(f"Error displaying chart preview: {e}")
        
        # Display CAGR chart if available
        if 'trend_analysis' in slide and slide['trend_analysis']:
            st.markdown("---")
            st.markdown("**Trend Analysis:**")
            
            # CAGR chart selection checkbox
            cagr_col1, cagr_col2 = st.columns([1, 3])
            with cagr_col1:
                st.session_state.selected_cagr_charts[f"{slide_id}_cagr"] = st.checkbox(
                    "Include CAGR Chart in Export",
                    value=st.session_state.selected_cagr_charts.get(f"{slide_id}_cagr", False),
                    key=f"cagr_select_{slide_id}",
                    help="Include the CAGR chart in the exported PowerPoint"
                )
            
            with cagr_col2:
                st.markdown("**CAGR Analysis**")
            
            # Display CAGR data
            try:
                trend_data = slide['trend_analysis']
                trend_rows = []
                
                for metric, values in trend_data.items():
                    # Skip moving average columns
                    if any(f"MA{p}" in metric for p in [2, 3, 5, 10]):
                        continue
                        
                    cagr = values.get('cagr', None)
                    recent_trend = values.get('recent_trend', None)
                    latest = values.get('latest', None)
                    
                    if cagr is not None and not pd.isna(cagr):
                        trend_rows.append({
                            "Metric": metric,
                            "Latest Value": f"${latest:,.2f}M" if pd.notna(latest) and isinstance(latest, (int, float)) else str(latest) if latest is not None else "N/A",
                            "CAGR (%)": float(f"{cagr:.2f}") if pd.notna(cagr) else None, # Store as float for plotting
                            "Recent Trend (%)": float(f"{recent_trend:.2f}") if pd.notna(recent_trend) else None, # Store as float
                        })
                
                if trend_rows:
                    trend_df = pd.DataFrame(trend_rows)
                    # st.dataframe(trend_df, use_container_width=True) # Removed raw data display
                    
                    # Create a bar chart for CAGR comparison
                    # Ensure 'CAGR (%)' column exists and is numeric
                    if 'CAGR (%)' in trend_df.columns and pd.api.types.is_numeric_dtype(trend_df['CAGR (%)']):
                        fig = px.bar(
                            trend_df.dropna(subset=['CAGR (%)']),
                            x="Metric", 
                            y="CAGR (%)",
                            title=f"{slide.get('title', '')} - Compound Annual Growth Rate (CAGR)",
                            color="CAGR (%)",
                            color_continuous_scale="RdYlGn",
                            height=400,
                            labels={"CAGR (%)": "CAGR Growth Rate (%)"}
                        )
                        fig.update_layout(
                            yaxis_title="CAGR (%)",
                            xaxis_title="Financial Metric",
                            margin=dict(l=20, r=20, t=40, b=20)
                        )
                        st.markdown("**Visualization Type:** Bar Chart for CAGR % by Metric")
                        st.plotly_chart(fig, use_container_width=True)
                    else:
                        st.info("No valid CAGR data to plot.")
                else:
                    st.info("No trend analysis data available or processed for display.")
                    
            except Exception as e:
                st.error(f"Error displaying trend analysis: {e}")

def main():
    st.set_page_config(
        page_title="Preview & Export",
        page_icon="👁️",
        layout="wide"
    )
    
    st.title("👁️ Preview & Export")
    
    if 'slides' not in st.session_state or not st.session_state.slides:
        st.warning("No slides to preview. Go to the Slide Builder to create some slides!")
        return
    
    st.markdown("""
    ### Preview your presentation
    Review your slides and select which ones to include in the exported PowerPoint.
    You can also choose to include or exclude individual charts and CAGR analyses.
    """)
    
    # Display slide previews
    for idx, slide in enumerate(st.session_state.slides):
        display_slide_preview(slide, idx)
    
    # Export controls
    st.markdown("---")
    st.subheader("Export Options")
    
    # Select all/none toggles
    col1, col2, col3 = st.columns(3)
    
    with col1:
        if st.button("Select All Slides", use_container_width=True):
            for slide in st.session_state.slides:
                slide_id = slide.get('id', f'slide_{st.session_state.slides.index(slide)}')
                st.session_state.selected_slides[slide_id] = True
        
        if st.button("Deselect All Slides", use_container_width=True):
            for slide in st.session_state.slides:
                slide_id = slide.get('id', f'slide_{st.session_state.slides.index(slide)}')
                st.session_state.selected_slides[slide_id] = False
    
    with col2:
        if st.button("Select All Charts", use_container_width=True):
            for slide in st.session_state.slides:
                slide_id = slide.get('id', f'slide_{st.session_state.slides.index(slide)}')
                if 'chart_data' in slide:
                    st.session_state.selected_charts[f"{slide_id}_chart"] = True
    
    with col3:
        if st.button("Select All CAGR Analyses", use_container_width=True):
            for slide in st.session_state.slides:
                slide_id = slide.get('id', f'slide_{st.session_state.slides.index(slide)}')
                if 'trend_analysis' in slide:
                    st.session_state.selected_cagr_charts[f"{slide_id}_cagr"] = True
    
    # Generate and download buttons
    st.markdown("---")
    st.subheader("Generate Presentation")
    
    col1, col2 = st.columns(2)
    
    with col1:
        if st.button("🔄 Generate PowerPoint", type="primary", use_container_width=True):
            with st.spinner("Creating your presentation..."):
                try:
                    pptx_bytes = create_presentation()
                    st.session_state.pptx_bytes = pptx_bytes
                    st.session_state.generation_time = datetime.now()
                    st.success("Presentation generated successfully!")
                except Exception as e:
                    st.error(f"Error generating presentation: {e}")
    
    with col2:
        if 'pptx_bytes' in st.session_state and st.session_state.pptx_bytes:
            # Create a download button with the current timestamp
            timestamp = st.session_state.generation_time.strftime("%Y%m%d_%H%M%S")
            st.download_button(
                label="💾 Download PowerPoint",
                data=st.session_state.pptx_bytes,
                file_name=f"financial_analysis_{timestamp}.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                use_container_width=True,
                type="primary"
            )
        else:
            st.button(
                "💾 Download PowerPoint", 
                disabled=True, 
                help="Generate the presentation first",
                use_container_width=True
            )
    
    # Add some styling
    st.markdown("""
    <style>
    .stButton>button {
        height: 3em;
    }
    .stButton>button[kind="primary"] {
        background-color: #4CAF50;
        color: white;
    }
    </style>
    """, unsafe_allow_html=True)

if __name__ == "__main__":
    main()
