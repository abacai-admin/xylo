import pandas as pd
import numpy as np
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_LEGEND_POSITION
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from typing import List, Dict, Any, Optional, Tuple

# Import chart colors utility
from .chart_colors import set_chart_colors

def create_table_slide(presentation: Presentation, title: str, data: pd.DataFrame, 
                       selected_columns: List[str]) -> None:
    """
    Create a slide with a table based on selected data
    
    Args:
        presentation: PowerPoint presentation object
        title: Title for the slide
        data: DataFrame containing the financial data
        selected_columns: List of column names to include in the table
    """
    # Create a blank slide (don't rely on layouts with placeholders)
    try:
        slide_layout = presentation.slide_layouts[1]
    except IndexError:
        slide_layout = presentation.slide_layouts[5] # Fallback to blank if 1 is not there
    slide = presentation.slides.add_slide(slide_layout)
    
    # Add title using placeholder if available, otherwise add textbox
    if slide.shapes.title:
        slide.shapes.title.text = title
    else:
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
        title_shape.text_frame.text = title
        title_para = title_shape.text_frame.paragraphs[0]
        title_para.font.size = Pt(24)
        title_para.font.bold = True
    
    # Debug information
    print(f"Creating table slide with title: {title}")
    print(f"Data columns: {data.columns.tolist()}")
    print(f"Selected columns: {selected_columns}")
    
    # Filter the DataFrame to include only selected columns
    if selected_columns:
        # Always include Year/Date column for reference
        if 'Year' in data.columns and 'Year' not in selected_columns:
            display_cols = ['Year'] + selected_columns
        elif 'Date' in data.columns and 'Date' not in selected_columns:
            display_cols = ['Date'] + selected_columns
        else:
            display_cols = selected_columns
            
        # Make sure all columns actually exist in the data
        display_cols = [col for col in display_cols if col in data.columns]
        
        # Filter the DataFrame
        display_data = data[display_cols]
    else:
        display_data = data
    
    # Check if we actually have data to display
    if display_data.empty:
        # Add message if no data
        msg_shape = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(1))
        msg_shape.text_frame.text = "No data available for the selected metrics"
        return
    
    # Define table dimensions
    rows, cols = len(display_data) + 1, len(display_data.columns)  # +1 for header row
    
    # Add a table to the slide with fixed position
    left = Inches(0.5)
    top = Inches(1.8)  # Below the title
    width = Inches(9)
    height = Inches(5)  # Taller to accommodate data
    
    # Add the table directly to the slide as a native PowerPoint table (fully editable)
    try:
        # Create a native PowerPoint table that can be edited after export
        table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
        table = table_shape.table
        
        # Set column headers
        for i, column_name in enumerate(display_data.columns):
            cell = table.cell(0, i)
            cell.text = str(column_name)  # Ensure it's a string
            
            # Format header cell - these formatting options will be editable in PowerPoint
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(92, 158, 173)  # Blue header
            
            # Format header text with standard PowerPoint text formatting
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.bold = True
            paragraph.font.size = Pt(12)
            paragraph.font.color.rgb = RGBColor(255, 255, 255)  # White text
        
        # Fill in data rows
        for row_idx, (_, row_data) in enumerate(display_data.iterrows(), start=1):
            for col_idx, value in enumerate(row_data):
                cell = table.cell(row_idx, col_idx)
                
                # Format numbers nicely
                if isinstance(value, (int, float)) and pd.notna(value):
                    # Format as currency with millions for financial metrics
                    if display_data.columns[col_idx] not in ['Year']:
                        cell.text = f"${value:,.2f}M"
                    else:
                        cell.text = str(value)
                else:
                    cell.text = str(value) if pd.notna(value) else ""
                
                # Format data cells with alternating colors
                if row_idx % 2 == 1:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(235, 241, 243)  # Light blue for odd rows
        
        # Set column widths based on content
        total_width = width.inches
        col_width = total_width / cols
        for col in table.columns:
            col.width = Inches(col_width)
            
        print(f"Table created successfully with {rows} rows and {cols} columns")
    except Exception as e:
        # If table creation fails, add an error message
        error_shape = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1))
        error_shape.text_frame.text = f"Error creating table: {str(e)}"
        print(f"Error creating table: {str(e)}")

def create_bar_chart_slide(presentation: Presentation, title: str, data: pd.DataFrame, 
                           selected_columns: List[str]) -> None:
    """
    Create a slide with a bar chart based on selected data
    
    Args:
        presentation: PowerPoint presentation object
        title: Title for the slide
        data: DataFrame containing the financial data
        selected_columns: List of column names to include in the bar chart
    """
    # Create a blank slide (don't rely on layouts with placeholders)
    try:
        slide_layout = presentation.slide_layouts[1]
    except IndexError:
        slide_layout = presentation.slide_layouts[5] # Fallback to blank
    slide = presentation.slides.add_slide(slide_layout)
    
    # Add a title manually or use placeholder
    if slide.shapes.title:
        slide.shapes.title.text = title
    else:
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
        title_shape.text_frame.text = title
        title_para = title_shape.text_frame.paragraphs[0]
        title_para.font.size = Pt(24)
        title_para.font.bold = True
    
    # Debug information
    print(f"Creating bar chart slide with title: {title}")
    print(f"Data columns: {data.columns.tolist()}")
    print(f"Selected columns: {selected_columns}")
    
    # Check if we have valid data and columns to plot
    if data.empty or not selected_columns:
        # Add a text box with an error message
        textbox = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(1))
        textbox.text_frame.text = "No data available for the selected metrics"
        return
    
    # Filter the DataFrame to include only selected columns
    # Make sure all columns actually exist in the data
    available_cols = [col for col in selected_columns if col in data.columns]
    
    if not available_cols:
        # Add a text box with an error message
        textbox = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(1))
        textbox.text_frame.text = "Selected metrics not found in the dataset"
        return
    
    try:
        # Get categories (years) for the chart
        if 'Year' in data.columns:
            categories = data['Year'].astype(str).tolist()
        elif 'Date' in data.columns:
            # Extract years from dates
            categories = pd.to_datetime(data['Date']).dt.year.astype(str).tolist()
        else:
            # If no Year/Date column, use row indices as categories
            categories = [f"Period {i+1}" for i in range(len(data))]
        
        # Create chart data
        chart_data = CategoryChartData()
        chart_data.categories = categories
        
        # Add series for each selected column
        for column in available_cols:
            # Skip non-numeric columns
            if pd.api.types.is_numeric_dtype(data[column]):
                # Extract values, replacing NaN with 0
                values = data[column].fillna(0).tolist()
                chart_data.add_series(column, values)
        
        # Define chart placement - centered on slide, below title
        x, y, cx, cy = Inches(1.0), Inches(1.8), Inches(8), Inches(5)
        
        # Create the chart directly on the slide
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
        ).chart
        
        # Format the chart
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM  # Proper enum value for bottom position
        chart.legend.include_in_layout = False
        
        # Format axes
        value_axis = chart.value_axis
        value_axis.has_major_gridlines = True
        value_axis.tick_labels.font.size = Pt(10)
        value_axis.has_title = True
        value_axis.axis_title.text_frame.text = "Amount (USD Millions)"
        
        category_axis = chart.category_axis
        category_axis.tick_labels.font.size = Pt(10)
        category_axis.has_title = True
        category_axis.axis_title.text_frame.text = "Year"
        
        # Add data labels
        plot = chart.plots[0]
        plot.has_data_labels = True
        data_labels = plot.data_labels
        data_labels.font.size = Pt(9)
        data_labels.position = 0  # Outside End
        
        # Apply custom brand colors to the chart
        set_chart_colors(chart)
        print(f"Applied custom brand colors to chart")
        
        print(f"Bar chart created successfully with {len(available_cols)} series")
    except Exception as e:
        # If chart creation fails, add an error message
        error_shape = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1))
        error_shape.text_frame.text = f"Error creating chart: {str(e)}"
        print(f"Error creating chart: {str(e)}")

def create_comparison_table_slide(presentation: Presentation, title: str, data: pd.DataFrame, 
                       selected_columns: List[str], ticker1: str, ticker2: str) -> None:
    """
    Create a slide with a comparison table for two companies
    
    Args:
        presentation: PowerPoint presentation object
        title: Title for the slide
        data: DataFrame containing the merged financial data
        selected_columns: List of base column names to include in the table
        ticker1: First company ticker
        ticker2: Second company ticker
    """
    # Create a blank slide (don't rely on layouts with placeholders)
    try:
        slide_layout = presentation.slide_layouts[1]
    except IndexError:
        slide_layout = presentation.slide_layouts[5] # Fallback to blank
    slide = presentation.slides.add_slide(slide_layout)
    
    # Add a title manually or use placeholder
    if slide.shapes.title:
        slide.shapes.title.text = title
    else:
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
        title_shape.text_frame.text = title
        title_para = title_shape.text_frame.paragraphs[0]
        title_para.font.size = Pt(24)
        title_para.font.bold = True
    
    # Debug information
    print(f"Creating comparison table slide with title: {title}")
    print(f"Data columns: {data.columns.tolist()}")
    print(f"Selected columns: {selected_columns}")
    
    # If no columns selected, try to identify suitable comparison metrics
    if not selected_columns:
        # Look for metrics that are available for both companies
        potential_metrics = []
        for col in data.columns:
            # Skip Year, Date columns
            if col in ['Year', 'Date']:
                continue
                
            # Extract the base metric name (remove company suffix)
            for suffix in [f"_{ticker1}", f"_{ticker2}"]:
                if col.endswith(suffix):
                    base_metric = col[:-len(suffix)]
                    # Check if this metric exists for both companies
                    if f"{base_metric}_{ticker1}" in data.columns and f"{base_metric}_{ticker2}" in data.columns:
                        potential_metrics.append(base_metric)
                    break
                    
        # Remove duplicates and limit to top 5 metrics
        selected_columns = list(dict.fromkeys(potential_metrics))[:5]
        
    # If we still don't have any columns, return an error
    if not selected_columns:
        textbox = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(1))
        textbox.text_frame.text = "No common metrics found for comparison"
        return
        
    # Check if we have Year column for comparison
    if 'Year' not in data.columns:
        textbox = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(1))
        textbox.text_frame.text = "Year data not available for comparison"
        return
        
    # Get all years with data
    years = sorted(data['Year'].unique())
    
    # Create table dimensions
    rows = len(years) + 1  # +1 for header row
    cols = len(selected_columns) * 2 + 1  # Each metric gets 2 columns (one per company) + 1 for years
    
    # Add a fully editable table to the slide
    left = Inches(0.5)
    top = Inches(2.0)
    width = Inches(9)
    height = Inches(4)
    
    # Create table with dimension for all metrics and both companies
    rows = len(years) + 1  # +1 for headers
    cols = 1 + len(selected_columns) * 2  # Year column + (metrics * 2 companies)
    
    # Create a native PowerPoint table that will remain fully editable after export
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    
    print(f"Creating editable comparison table with {rows} rows and {cols} columns")

    # Set first column header (Year)
    cell = table.cell(0, 0)
    cell.text = "Year"
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(92, 158, 173)  # Blue header
    
    # Format header text
    paragraph = cell.text_frame.paragraphs[0]
    paragraph.font.bold = True
    paragraph.font.size = Pt(12)
    paragraph.font.color.rgb = RGBColor(255, 255, 255)  # White text
    
    # Set metric column headers
    col_idx = 1
    for metric in selected_columns:
        # Header for first company
        cell = table.cell(0, col_idx)
        cell.text = f"{metric} - {ticker1}"
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(92, 158, 173)  # Blue header
        
        # Format header text
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.bold = True
        paragraph.font.size = Pt(12)
        paragraph.font.color.rgb = RGBColor(255, 255, 255)  # White text
        
        # Header for second company
        cell = table.cell(0, col_idx + 1)
        cell.text = f"{metric} - {ticker2}"
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(92, 158, 173)  # Blue header
        
        # Format header text
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.bold = True
        paragraph.font.size = Pt(12)
        paragraph.font.color.rgb = RGBColor(255, 255, 255)  # White text
        
        col_idx += 2
    
    # Fill in data rows
    for row_idx, year in enumerate(years, start=1):
        # Set year in first column
        cell = table.cell(row_idx, 0)
        cell.text = str(int(year)) if isinstance(year, (int, float)) else str(year)
        
        # Add alternating row colors
        if row_idx % 2 == 1:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(235, 241, 243)  # Light blue for odd rows
        
        # Fill in metric values for both companies
        col_idx = 1
        for metric in selected_columns:
            # Get the data for this year
            year_data = data[data['Year'] == year]
            
            # Fill in first company value
            cell = table.cell(row_idx, col_idx)
            col1_name = f"{metric}_{ticker1}"
            if not year_data.empty and col1_name in year_data.columns:
                value = year_data[col1_name].iloc[0] if not year_data[col1_name].isna().all() else None
                if pd.notna(value):
                    cell.text = f"${value:,.2f}M"
                else:
                    cell.text = "N/A"
            else:
                cell.text = "N/A"
            
            # Add alternating row colors
            if row_idx % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(235, 241, 243)  # Light blue for odd rows
            
            # Fill in second company value
            cell = table.cell(row_idx, col_idx + 1)
            col2_name = f"{metric}_{ticker2}"
            if not year_data.empty and col2_name in year_data.columns:
                value = year_data[col2_name].iloc[0] if not year_data[col2_name].isna().all() else None
                if pd.notna(value):
                    cell.text = f"${value:,.2f}M"
                else:
                    cell.text = "N/A"
            else:
                cell.text = "N/A"
                
            # Add alternating row colors
            if row_idx % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(235, 241, 243)  # Light blue for odd rows
                
            col_idx += 2
    
    # Autofit columns
    for col in table.columns:
        col.width = Inches(width.inches / cols)

def create_comparison_bar_chart_slide(presentation: Presentation, title: str, data: pd.DataFrame,
                                     selected_metrics: List[str], ticker1: str, ticker2: str) -> None:
    """
    Create a slide with comparison bar charts for two companies
    
    Args:
        presentation: PowerPoint presentation object
        title: Title for the slide
        data: DataFrame containing the merged financial data
        selected_metrics: List of metrics to compare
        ticker1: First company ticker
        ticker2: Second company ticker
    """
    # For each selected metric, create a separate chart slide
    for metric in selected_metrics[:3]:  # Limit to 3 metrics to avoid too many slides
        # Create a slide using a standard layout
        try:
            slide_layout = presentation.slide_layouts[1] # Title and Content
        except IndexError:
            slide_layout = presentation.slide_layouts[5] # Fallback to Blank
        slide = presentation.slides.add_slide(slide_layout)
        
        # Add a title manually or use placeholder
        metric_title = f"{title} - {metric} Comparison"
        if slide.shapes.title:
            slide.shapes.title.text = metric_title
        else:
            title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
            title_shape.text_frame.text = metric_title
            title_para = title_shape.text_frame.paragraphs[0]
            title_para.font.size = Pt(24)
            title_para.font.bold = True
        
        print(f"Creating comparison chart for metric: {metric}")
        
        # Get years and prepare chart data
        if 'Year' not in data.columns:
            # Without year column, we can't create a proper chart
            textbox = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(1))
            textbox.text_frame.text = "No year data available for comparison"
            continue
            
        years = sorted(data['Year'].unique())
        
        # Create chart data for this metric
        chart_data = CategoryChartData()
        chart_data.categories = [str(int(year)) if isinstance(year, (int, float)) else str(year) for year in years]
        
        # Extract values for both companies
        col1_name = f"{metric}_{ticker1}"
        col2_name = f"{metric}_{ticker2}"
        
        # Check if the columns exist
        values1 = []
        values2 = []
        
        for year in years:
            year_data = data[data['Year'] == year]
            
            # Get value for first company
            if not year_data.empty and col1_name in year_data.columns:
                val1 = year_data[col1_name].iloc[0] if not year_data[col1_name].isna().all() else 0
                values1.append(float(val1) if pd.notna(val1) else 0)
            else:
                values1.append(0)
                
            # Get value for second company
            if not year_data.empty and col2_name in year_data.columns:
                val2 = year_data[col2_name].iloc[0] if not year_data[col2_name].isna().all() else 0
                values2.append(float(val2) if pd.notna(val2) else 0)
            else:
                values2.append(0)
        
        # Add data series for both companies
        chart_data.add_series(ticker1, values1)
        chart_data.add_series(ticker2, values2)
        
        # Define chart placement - centered on slide, below title
        x, y, cx, cy = Inches(1.0), Inches(1.8), Inches(8), Inches(5)
        
        # Create the chart directly on the slide - using native PowerPoint chart objects so they're fully editable
        try:
            # Use native PowerPoint chart type to ensure editability
            chart_shape = slide.shapes.add_chart(
                XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
            )
            chart = chart_shape.chart
            
            # Make sure data is properly linked for editability
            # PowerPoint automatically creates an Excel worksheet behind the scenes
            # that stores the chart data and allows for future editing
            
            print(f"Successfully added editable chart for {metric}")
        except Exception as e:
            # If chart creation fails, add an error message
            error_shape = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1))
            error_shape.text_frame.text = f"Error creating chart: {str(e)}"
            print(f"Error creating chart for {metric}: {str(e)}")
            continue
        
        # Format the chart
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM  # Proper enum value for bottom position
        chart.legend.include_in_layout = False
        
        # Format axes
        value_axis = chart.value_axis
        value_axis.has_major_gridlines = True
        value_axis.tick_labels.font.size = Pt(10)
        value_axis.has_title = True
        value_axis.axis_title.text_frame.text = f"{metric} (USD Millions)"
        
        category_axis = chart.category_axis
        category_axis.tick_labels.font.size = Pt(10)
        category_axis.has_title = True
        category_axis.axis_title.text_frame.text = "Year"
        
        # Add data labels
        plot = chart.plots[0]
        plot.has_data_labels = True
        data_labels = plot.data_labels
        data_labels.font.size = Pt(9)
        data_labels.position = 0  # Outside End
        
        # Apply custom brand colors to the comparison chart
        set_chart_colors(chart)
        print(f"Applied custom brand colors to comparison chart for {metric}")

def create_cagr_chart_slide(presentation: Presentation, title: str, trend_data: Dict[str, Any]) -> None:
    """Create a slide that shows CAGR bar chart based on trend_analysis dict."""
    # We expect trend_data to be a dict {metric: {cagr: value, ...}}
    # Collect metrics and CAGR values
    metrics = []
    cagr_vals = []
    for metric, vals in trend_data.items():
        if any(f"MA{p}" in metric for p in [2, 3, 5, 10]):
            continue  # skip moving average synthetic metrics
        cagr_val = vals.get("cagr")
        if cagr_val is not None and pd.notna(cagr_val):
            metrics.append(metric)
            cagr_vals.append(float(cagr_val))

    if not metrics:
        # nothing to chart – skip slide creation
        return

    # Pick a slide layout (title+content or fallback)
    try:
        slide_layout = presentation.slide_layouts[1]
    except IndexError:
        slide_layout = presentation.slide_layouts[0]

    slide = presentation.slides.add_slide(slide_layout)

    # Title
    if slide.shapes.title:
        slide.shapes.title.text = title
    else:
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
        title_box.text_frame.text = title

    # Prepare chart data
    chart_data = CategoryChartData()
    chart_data.categories = metrics
    chart_data.add_series("CAGR %", cagr_vals)

    # Place chart
    x, y, cx, cy = Inches(1), Inches(1.8), Inches(8), Inches(5)
    chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data).chart

    # Format chart
    chart.has_legend = False
    value_axis = chart.value_axis
    value_axis.has_major_gridlines = True
    value_axis.tick_labels.font.size = Pt(10)
    value_axis.has_title = True
    value_axis.axis_title.text_frame.text = "CAGR (%)"
    category_axis = chart.category_axis
    category_axis.tick_labels.font.size = Pt(10)

    # Add data labels
    plot = chart.plots[0]
    plot.has_data_labels = True
    plot.data_labels.number_format = "0.0%"

    # Apply brand colors
    set_chart_colors(chart)

def get_template_path():
    """Get the path to the template file"""
    import os
    # Get the directory of the current file
    current_dir = os.path.dirname(os.path.abspath(__file__))
    # Go up one level to the project root, then into assets
    return os.path.join(current_dir, '..', 'assets', 'template.pptx')

def generate_presentation(slides_data: List[Dict[str, Any]], output_path: str) -> str:
    template_path = get_template_path()
    prs = None
    title_slide_layout = None
    content_slide_layout = None
    closing_slide_layout = None

    try:
        # ------------------------------------------------------------------
        # Create the presentation based on the template and keep the
        # original opening (slide 1) and closing (last slide). Remove
        # any slides that sit in between so our generated content can
        # be inserted while the opening and closing slides keep their
        # original styling.
        # ------------------------------------------------------------------

        prs = Presentation(template_path)

        # Keep ALL template slides intact. We'll later duplicate Slide 2 for
        # each generated chart/table, but we no longer remove any slides so
        # the original opening (slide 1) and closing (last slide) remain
        # untouched.

        preserved_template_closing = len(prs.slides) > 1  # True if template has a closing slide
        if preserved_template_closing:
            original_closing_slide_id = prs.slides[-1].slide_id  # store id before generating new slides

        # Re-assign layouts after we recreated `prs`
        if len(prs.slide_layouts) > 0:
            title_slide_layout = prs.slide_layouts[0]
        if len(prs.slide_layouts) > 1:
            content_slide_layout = prs.slide_layouts[1]
        if len(prs.slide_layouts) > 5:
            closing_slide_layout = prs.slide_layouts[5]
        else:
            closing_slide_layout = prs.slide_layouts[0]

    except Exception as e:
        print(f"Warning: Could not load or process template from '{template_path}': {e}. Falling back to default presentation.")
        prs = Presentation()
        prs.core_properties.title = "Financial Data Presentation"
        prs.core_properties.author = "Presentation Builder"

        # Use default layouts for the new presentation
        title_slide_layout = prs.slide_layouts[0]
        content_slide_layout = prs.slide_layouts[1]
        closing_slide_layout = prs.slide_layouts[0] # Use title slide layout for closing if no specific one

        current_title_slide = prs.slides.add_slide(title_slide_layout)
        if current_title_slide.shapes.title: current_title_slide.shapes.title.text = "Financial Data Presentation"
        if len(current_title_slide.placeholders) > 1: current_title_slide.placeholders[1].text_frame.text = f"Generated on {pd.Timestamp.now().strftime('%Y-%m-%d')}"

    # Ensure content_slide_layout is set (it should be by now)
    if content_slide_layout is None:
        content_slide_layout = prs.slide_layouts[1] # Fallback to default title and content

    # Debug information
    print(f"\nGenerating presentation with {len(slides_data)} slides using content layout: {content_slide_layout.name if content_slide_layout else 'Default'}")
    for i, slide_data in enumerate(slides_data):
        print(f"\nSlide {i+1} details:")
        print(f"  Title: {slide_data.get('title')}")
        print(f"  Ticker: {slide_data.get('ticker')}")
    
    # Process each slide
    for slide_item_data in slides_data:
        try:
            if 'chart_data' not in slide_item_data or slide_item_data.get('chart_data') is None:
                print("Skipping slide: No chart data found")
                continue
            data = slide_item_data.get('chart_data')
            if not isinstance(data, pd.DataFrame) or data.empty:
                print("Skipping slide: Invalid or empty chart data")
                continue
            selected_columns = [col for col in slide_item_data.get('selected_columns', []) if col in data.columns]
            if not selected_columns:
                selected_columns = data.select_dtypes(include=np.number).columns.tolist()
                if not selected_columns:
                    print("Skipping slide: No valid numeric columns found for chart after attempting fallback")
                    continue
                print(f"Warning: No specific columns selected or found. Using all numeric columns: {selected_columns}")

            title = slide_item_data.get('title', 'Financial Data')
            ticker = slide_item_data.get('ticker', '')
            ticker2 = slide_item_data.get('ticker2', '')
            chart_type = slide_item_data.get('chart_type', 'table')
            is_comparison = slide_item_data.get('is_comparison', False)
            formatted_title = f"{title} - {ticker} vs {ticker2}" if is_comparison and ticker2 else f"{title} - {ticker}" if ticker else title

            # Call helper functions (they will add slides to prs)
            if is_comparison and 'chart_data2' in slide_item_data and isinstance(slide_item_data['chart_data2'], pd.DataFrame) and not slide_item_data['chart_data2'].empty:
                data2 = slide_item_data['chart_data2']
                data1_renamed = data.rename(columns={col: f"{col}_{ticker}" for col in data.columns if col not in ['Year', 'Date']})
                data2_renamed = data2.rename(columns={col: f"{col}_{ticker2}" for col in data2.columns if col not in ['Year', 'Date']})
                if 'Year' not in data1_renamed.columns and 'Date' in data1_renamed.columns:
                    data1_renamed['Year'] = pd.to_datetime(data1_renamed['Date']).dt.year
                if 'Year' not in data2_renamed.columns and 'Date' in data2_renamed.columns:
                    data2_renamed['Year'] = pd.to_datetime(data2_renamed['Date']).dt.year
                
                if 'Year' in data1_renamed.columns and 'Year' in data2_renamed.columns:
                    merged_data = pd.merge(data1_renamed, data2_renamed, on='Year', how='outer')
                    base_metrics_for_comparison = selected_columns
                    if chart_type == "bar_chart":
                        create_comparison_bar_chart_slide(prs, formatted_title, merged_data, base_metrics_for_comparison, ticker, ticker2)
                    else:
                        create_comparison_table_slide(prs, formatted_title, merged_data, base_metrics_for_comparison, ticker, ticker2)
                else:
                    print(f"Skipping comparison for '{formatted_title}' due to missing 'Year' column for merging. Creating slide for first ticker only.")
                    data_cols_to_use = [col for col in selected_columns if col in data.columns] or data.select_dtypes(include=np.number).columns.tolist()
                    if chart_type == "bar_chart":
                        create_bar_chart_slide(prs, f"{title} - {ticker}", data, data_cols_to_use)
                    else:
                        create_table_slide(prs, f"{title} - {ticker}", data, data_cols_to_use)
            else:
                data_cols_to_use = [col for col in selected_columns if col in data.columns] or data.select_dtypes(include=np.number).columns.tolist()
                if chart_type == "bar_chart":
                    create_bar_chart_slide(prs, formatted_title, data, data_cols_to_use)
                else:
                    create_table_slide(prs, formatted_title, data, data_cols_to_use)

            # If trend analysis (CAGR) exists – add a CAGR chart slide
            if slide_item_data.get("trend_analysis"):
                try:
                    create_cagr_chart_slide(prs, f"{title} - CAGR Analysis", slide_item_data["trend_analysis"])
                except Exception as e_cagr:
                    print(f"Warning: Could not create CAGR slide for '{title}': {e_cagr}")

        except Exception as e_outer_slide_processing:
            print(f"Outer error processing slide '{slide_item_data.get('title', 'Unknown')}': {e_outer_slide_processing}")
            try:
                error_fallback_slide = prs.slides.add_slide(prs.slide_layouts[5]) # Blank layout
                textbox = error_fallback_slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(5))
                textbox.text_frame.text = f"Error processing slide: {slide_item_data.get('title', 'Unknown')}\nDetails: {e_outer_slide_processing}"
                print("Added error notification slide.")
            except Exception as e_error_slide:
                print(f"Could not even add an error notification slide: {e_error_slide}")
    
    # ------------------------------------------------------------------
    # Handle the closing slide.
    # If we preserved the template's closing slide, we need to move it to
    # the end because newly-generated content slides were appended after it.
    # If no closing slide existed in the template we add a generic one.
    # ------------------------------------------------------------------
    if 'preserved_template_closing' in locals() and preserved_template_closing:
        try:
            # Locate the original closing slide in the (possibly reordered) list
            sldIdLst = prs.slides._sldIdLst
            closing_idx = None
            for idx, slide in enumerate(prs.slides):
                if slide.slide_id == original_closing_slide_id:
                    closing_idx = idx
                    break

            # Move only if it is not already last
            if closing_idx is not None and closing_idx != len(prs.slides) - 1:
                closing_elem = sldIdLst[closing_idx]
                sldIdLst.remove(closing_elem)
                sldIdLst.append(closing_elem)
        except Exception as e_move:
            print(f"Warning: Could not reposition closing slide: {e_move}")
    else:
        # No closing slide in template – add a simple one using the chosen layout
        try:
            final_closing_slide = prs.slides.add_slide(closing_slide_layout if closing_slide_layout else prs.slide_layouts[0])
            try:
                if final_closing_slide.shapes.title:
                    final_closing_slide.shapes.title.text = "Thank You"
                if len(final_closing_slide.placeholders) > 1 and final_closing_slide.placeholders[1].has_text_frame:
                    final_closing_slide.placeholders[1].text_frame.text = "Thank you for your attention"
                elif not final_closing_slide.shapes.title:
                    # Fallback textbox for truly blank layouts
                    title_shape = final_closing_slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2))
                    title_shape.text_frame.text = "Thank You"
                    for para in title_shape.text_frame.paragraphs:
                        para.font.size = Pt(44)
                        para.font.bold = True
                        para.alignment = PP_ALIGN.CENTER
            except Exception as e_populate:
                print(f"Warning: Could not fully populate closing slide: {e_populate}")
        except Exception as e_add_closing:
            print(f"Warning: Could not add closing slide: {e_add_closing}")
    
    # Save the presentation
    try:
        prs.save(output_path)
        return output_path
    except Exception as e:
        print(f"Error saving presentation: {e}")
        raise
    finally:
        # Ensure the function always returns a value or raises an exception
        # This block can be used for cleanup if needed, but in this case,
        # we're just ensuring the try-except structure is complete.
        pass
